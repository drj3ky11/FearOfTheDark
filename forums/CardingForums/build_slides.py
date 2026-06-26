#!/usr/bin/env python3
"""
Build PPTX for "Análisis Forense de Foros Underground — Carding Forums".

Run with:
    uv run python talks/01_carding_forums/build_slides.py

When you have a branded template, pass it as the first argument:
    uv run python talks/01_carding_forums/build_slides.py template.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── Color palette ────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy — main background
C_BG_SEC  = RGBColor(0x0F, 0x0F, 0x20)  # deeper navy — section dividers
C_ACCENT  = RGBColor(0xE9, 0x4E, 0x4E)  # red — section titles, highlights
C_TITLE   = RGBColor(0xFF, 0xFF, 0xFF)  # white — slide titles
C_BODY    = RGBColor(0xB8, 0xC0, 0xCC)  # blue-gray — body text
C_MUTED   = RGBColor(0x60, 0x68, 0x78)  # muted — secondary labels
C_CODE_BG = RGBColor(0x0D, 0x1B, 0x2A)  # very dark — code boxes
C_BADGE   = RGBColor(0x2A, 0x2A, 0x50)  # mid-navy — tag/badge background

# ─── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ─── Low-level helpers ────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _add_rect(slide, x, y, w, h, fill: RGBColor, alpha: int = 255):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()  # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def _add_txbox(
    slide,
    text: str,
    x, y, w, h,
    size: int = 20,
    color: RGBColor = C_BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    """Add a single-paragraph text box."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def _add_multiline(
    slide,
    lines: list[tuple],  # (text, size, color, bold, align)
    x, y, w, h,
) -> None:
    """Add a text box with multiple paragraphs, each with independent style."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True

    for i, (text, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"


def _add_bullets(
    slide,
    title: str,
    bullets: list[str | tuple],
    x=Inches(0.6), y=Inches(1.5),
    w=Inches(12.1), h=Inches(5.5),
    title_y=Inches(0.35),
    title_size=32,
    bullet_size=20,
) -> None:
    """
    Add title + bulleted list.
    bullets can be str (normal) or ("text", indent_level) for sub-bullets.
    """
    _add_txbox(slide, title, Inches(0.6), title_y, Inches(12.1), Inches(0.9),
               size=title_size, color=C_TITLE, bold=True)

    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT

        # Indent sub-bullets
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        indent = Emu(Inches(0.3 * level))
        pPr.set("marL", str(int(indent)))
        pPr.set("indent", str(int(-Inches(0.25))))

        run = p.add_run()
        prefix = "▸  " if level == 0 else "·  "
        run.text = prefix + text
        run.font.size = Pt(bullet_size if level == 0 else bullet_size - 3)
        run.font.color.rgb = C_BODY if level == 0 else C_MUTED
        run.font.name = "Calibri"
        run.font.bold = False


# ─── Slide builders ───────────────────────────────────────────────────────────

def title_slide(prs: Presentation, title: str, subtitle: str, tag: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, C_BG)

    # Left accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, C_ACCENT)

    # Title
    _add_txbox(slide, title,
               Inches(0.6), Inches(2.0), Inches(12.0), Inches(1.8),
               size=44, color=C_TITLE, bold=True)

    # Subtitle
    _add_txbox(slide, subtitle,
               Inches(0.6), Inches(3.9), Inches(12.0), Inches(1.0),
               size=24, color=C_BODY)

    # Tag badge
    _add_rect(slide, Inches(0.6), Inches(5.2), Inches(3.6), Inches(0.45), C_BADGE)
    _add_txbox(slide, tag,
               Inches(0.7), Inches(5.22), Inches(3.5), Inches(0.4),
               size=14, color=C_MUTED)

    # Bottom accent line
    _add_rect(slide, Inches(0.6), Inches(6.9), Inches(11.5), Inches(0.05), C_ACCENT)


def section_slide(prs: Presentation, number: str, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_SEC)

    # Full-width accent bar at top
    _add_rect(slide, Inches(0), Inches(0), W, Inches(0.1), C_ACCENT)

    # Large section number (background watermark feel)
    _add_txbox(slide, number,
               Inches(9.5), Inches(0.8), Inches(3.5), Inches(5.0),
               size=200, color=RGBColor(0x25, 0x25, 0x45), bold=True, align=PP_ALIGN.RIGHT)

    # Section label
    _add_txbox(slide, "SECCIÓN",
               Inches(0.6), Inches(2.2), Inches(6.0), Inches(0.5),
               size=13, color=C_ACCENT, bold=True)

    # Title
    _add_txbox(slide, title,
               Inches(0.6), Inches(2.7), Inches(9.0), Inches(1.8),
               size=40, color=C_TITLE, bold=True)

    # Subtitle
    if subtitle:
        _add_txbox(slide, subtitle,
                   Inches(0.6), Inches(4.6), Inches(9.0), Inches(0.8),
                   size=20, color=C_BODY)


def content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str | tuple],
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG)
    _add_rect(slide, Inches(0.6), Inches(1.2), Inches(11.5), Inches(0.03), C_ACCENT)
    _add_bullets(slide, title, bullets)
    if note:
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.4),
                   size=12, color=C_MUTED)


def two_col_slide(
    prs: Presentation,
    title: str,
    left_title: str, left_items: list[str],
    right_title: str, right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG)
    _add_rect(slide, Inches(0.6), Inches(1.2), Inches(11.5), Inches(0.03), C_ACCENT)
    _add_txbox(slide, title, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9),
               size=32, color=C_TITLE, bold=True)

    col_w = Inches(5.6)
    col_h = Inches(5.0)

    for col_x, col_title, items in [
        (Inches(0.6),  left_title,  left_items),
        (Inches(7.1), right_title, right_items),
    ]:
        _add_txbox(slide, col_title, col_x, Inches(1.5), col_w, Inches(0.5),
                   size=18, color=C_ACCENT, bold=True)

        txb = slide.shapes.add_textbox(col_x, Inches(2.1), col_w, col_h)
        tf = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = "▸  " + item
            run.font.size = Pt(18)
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"


def table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG)
    _add_rect(slide, Inches(0.6), Inches(1.2), Inches(11.5), Inches(0.03), C_ACCENT)
    _add_txbox(slide, title, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9),
               size=32, color=C_TITLE, bold=True)

    cols = len(headers)
    col_w = Inches(11.5 / cols)
    row_h = Inches(0.48)
    table_top = Inches(1.45)

    # Header row
    _add_rect(slide, Inches(0.6), table_top, Inches(11.5), row_h, C_ACCENT)
    for j, h in enumerate(headers):
        _add_txbox(slide, h,
                   Inches(0.65) + col_w * j, table_top + Emu(60000),
                   col_w - Inches(0.1), row_h,
                   size=14, color=C_TITLE, bold=True)

    # Data rows
    for i, row in enumerate(rows):
        row_y = table_top + row_h * (i + 1)
        bg = RGBColor(0x22, 0x22, 0x40) if i % 2 == 0 else RGBColor(0x1E, 0x1E, 0x38)
        _add_rect(slide, Inches(0.6), row_y, Inches(11.5), row_h, bg)
        for j, cell in enumerate(row):
            _add_txbox(slide, cell,
                       Inches(0.65) + col_w * j, row_y + Emu(50000),
                       col_w - Inches(0.1), row_h,
                       size=13, color=C_BODY)

    if note:
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.4),
                   size=12, color=C_MUTED)


def code_slide(
    prs: Presentation,
    title: str,
    description: str,
    code_lines: list[str],
    note: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG)
    _add_rect(slide, Inches(0.6), Inches(1.2), Inches(11.5), Inches(0.03), C_ACCENT)
    _add_txbox(slide, title, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9),
               size=32, color=C_TITLE, bold=True)

    if description:
        _add_txbox(slide, description, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.6),
                   size=18, color=C_BODY)

    desc_h = Inches(0.7) if description else Inches(0)
    code_top = Inches(1.9) + desc_h
    code_h = Inches(4.8) - desc_h

    _add_rect(slide, Inches(0.5), code_top, Inches(12.3), code_h, C_CODE_BG)

    code_text = "\n".join(code_lines)
    txb = slide.shapes.add_textbox(Inches(0.7), code_top + Inches(0.15),
                                   Inches(12.0), code_h - Inches(0.3))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)  # green — monospace code feel
    run.font.name = "Consolas"

    if note:
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.4),
                   size=12, color=C_MUTED)


def demo_slide(prs: Presentation, section_title: str, steps: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_SEC)
    _add_rect(slide, Inches(0), Inches(0), W, Inches(0.1), C_ACCENT)

    _add_txbox(slide, "DEMO EN VIVO",
               Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.7),
               size=14, color=C_ACCENT, bold=True)
    _add_txbox(slide, section_title,
               Inches(0.6), Inches(1.65), Inches(12.0), Inches(1.2),
               size=36, color=C_TITLE, bold=True)

    _add_txbox(slide, "notebooks/cases/01_carding_forums.ipynb",
               Inches(0.6), Inches(2.9), Inches(8.0), Inches(0.5),
               size=16, color=C_MUTED)

    txb = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(12.0), Inches(3.0))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{i+1}.  {step}"
        run.font.size = Pt(19)
        run.font.color.rgb = C_BODY
        run.font.name = "Calibri"


# ─── Slide content ────────────────────────────────────────────────────────────

def build(prs: Presentation) -> None:

    # ── Title ──────────────────────────────────────────────────────────────────
    title_slide(prs,
        "Análisis Forense de Foros Underground",
        "Del leak al perfil: carding forums como caso de estudio",
        "CSBC26  ·  OpHarvestSeason  ·  2009–2021",
    )

    # ══════════════════════════════════════════════════════════════════════════
    # SECCIÓN 1: EL ECOSISTEMA
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "01", "El ecosistema",
                  "Qué es carding, cómo funcionan estos foros y por qué terminan leakeados")

    content_slide(prs, "¿Qué es carding?", [
        "Uso y tráfico de datos de tarjetas robados (dumps, CVVs, fullz)",
        "Fullz = paquete de identidad completa: nombre, DNI, SSN, tarjeta, dirección",
        "Servicios asociados: checkers, cashers, droppers, coders, money mules",
        "Mercado estimado: >$32B/año en pérdidas globales (Nilson Report 2022)",
        ("Los foros son la infraestructura social: reputación, feedback, vouch", 1),
        ("Sin reputación no hay ventas — igual que cualquier marketplace", 1),
    ])

    content_slide(prs, "Los foros como marketplace criminal", [
        "Misma lógica que eBay o MercadoLibre, pero para datos robados",
        "Vendedor publica: tipo de dato, banco emisor, país, precio",
        "Comprador revisa reputación, paga en criptomonedas (BTC, XMR)",
        "Sistema de vouching: otros miembros confirman que el vendedor es legítimo",
        "Moderadores resuelven disputas (escrow)",
        "Alta especialización: algunos solo venden, otros solo compran, otros dan soporte técnico",
    ])

    content_slide(prs, "Timeline: evolución 2007–2021", [
        "2007–2010  Primera generación: Carder.su, Carding.biz — comunidades pequeñas, rusas",
        "2010–2015  Consolidación: Carders.cc, Cardersplanet, Carder.pro — crecimiento masivo",
        "2015–2018  Madurez: CardingMafia, Crdshop, Elitecarders — especialización por región",
        "2019–2021  Presión de LE + leaks internos → fragmentación y migración a Telegram",
        "2022+      Telegram como nuevo canal dominante, foros como archivo histórico",
    ], note="Los 10 foros que analizamos cubren toda esta evolución")

    content_slide(prs, "¿Por qué terminan leakeados?", [
        "Rival groups — otra organización roba la BD y la publica como golpe de imagen",
        "OpSec failures — admin expone la BD por misconfiguration (S3 público, backup sin auth)",
        "Law enforcement — incautación de servidores; datos publicados post-investigación",
        "Exit scam — el admin vende la BD antes de cerrar el foro y desaparecer con el dinero",
        "Disgruntled insider — moderador o admin descontento filtra como venganza",
        ("En todos los casos: el atacante ya tenía acceso — la pregunta es POR QUÉ lo tenía", 1),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SECCIÓN 2: LOS DATOS
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "02", "Los datos",
                  "OpHarvestSeason: qué hay en los dumps y qué información de valor contienen")

    content_slide(prs, "OpHarvestSeason", [
        "Colección de dumps de foros underground compilada para análisis forense",
        "12 GB comprimidos — cientos de millones de registros",
        "Categorías: Carding, Ransomware, DDoS, Hacking, Drugs, Far Right, Doxxing",
        "Este caso de estudio: Carding Forums — 10 dumps, 2009–2021",
        "Total: 1.13M usuarios · 1.52M posts · 10 foros",
        ("Los datos son sensibles — todo el procesamiento es local, sin APIs externas", 1),
    ])

    table_slide(prs, "Los 10 foros de carding", [
        "Foro", "Leak", "Formato", "Encoding", "Notas",
    ], [
        ["Carder.su",          "2009-02", "vBulletin SQL", "cp1251", "Solo users, 0 posts (dump parcial)"],
        ["Carding.biz",        "2009-11", "vBulletin SQL", "cp1251", "Primera gen"],
        ["Carders.cc",         "2010-12", "vBulletin SQL", "cp1251", ""],
        ["Cardersplanet.biz",  "2010-05", "vBulletin SQL", "cp1251", ""],
        ["Carder.pro",         "2013-04", "vBulletin SQL", "cp1251", "Ruso, prefijo vb_"],
        ["Cardingmafia.ws",    "2016-02", "Flat file :sep", "UTF-8",  "Sin SQL — delimitado por ':'"],
        ["Crdshop.su",         "2016-11", "vBulletin SQL", "UTF-8",  "INSERT con columnas explícitas"],
        ["Elitecarders.name",  "2016-08", "vBulletin SQL", "cp1251", ""],
        ["CardingMafia",       "2021-03", "vBulletin SQL", "cp1251", "Posible duplicado de Cardmafia.cc"],
        ["Cardmafia.cc",       "2021-03", "vBulletin SQL", "cp1251", "175MB — mismo tamaño que el anterior"],
    ], note="Mismo schema vBulletin en todos → un solo parser los maneja (con variantes)")

    content_slide(prs, "Formato: vBulletin SQL", [
        "vBulletin: plataforma de foros dominante en los 2000s–2010s (PHP + MySQL)",
        "Los dumps son exports estándar de mysqldump — INSERT INTO por tabla",
        "Encoding cp1251 (Cirílico Windows) en la mayoría — artefacto de la época",
        "Schema consistente: user, post, pmtext, userfield, thread, forum",
        "5 variantes encontradas en este dataset — el parser las maneja todas",
    ])

    code_slide(prs, "Las 5 variantes de formato", "El parser tiene que manejar todas sin cambiar el código de análisis", [
        "-- Variante 1: Standard vBulletin (mayoría)",
        "INSERT INTO user VALUES (1,'admin','hash','admin@mail.ru','Moscow',3,...);",
        "",
        "-- Variante 2: Prefijo vb_ en tablas (Carder.pro)",
        "INSERT INTO vb_user VALUES (...);",
        "",
        "-- Variante 3: INSERT con columnas explícitas + tabs (CardingMafia)",
        "INSERT INTO `user` (`userid`,\t`username`,\t`password`,...) VALUES (...);",
        "",
        "-- Variante 4: UTF-8 en lugar de cp1251 (Crdshop.su)",
        "-- (solo encoding cambia, el resto igual)",
        "",
        "-- Variante 5: Flat file sin SQL (Cardingmafia.ws)",
        "username:email:password:icq:joindate",
    ], note="Auto-detect: peek 2KB → UTF-8 first, fallback cp1251; detectar prefijo vb_ en runtime")

    content_slide(prs, "¿Qué hay en un dump?", [
        "user — el registro central: userid, username, email, password hash, ICQ, Skype, IP, joindate, posts, timezone",
        "post — el contenido: userid, dateline (timestamp), pagetext (texto del post), visible",
        "pmtext — mensajes privados entre usuarios (muy sensibles, muy reveladores)",
        "userfield — campos extra: Jabber, dirección Bitcoin, bio libre",
        ("La tabla user es suficiente para correlación cross-foro", 1),
        ("Los posts son necesarios para timezone inference y estilometría", 1),
        ("Los PMs tienen conversaciones explícitas sobre transacciones", 1),
    ])

    content_slide(prs, "Datos de valor para atribución", [
        "ICQ number — identificador duro, raramente cambiado, muy fiable para cross-forum",
        "Email — throwaway pero consistente si el usuario descuida la opsec",
        "Password hash — mismo hash en foros distintos = misma contraseña = misma persona",
        "IP address — valiosa si no usaba VPN consistentemente (muchos no lo hacían)",
        "Username — alta colisión en nombres cortos; gold en nicknames únicos y largos",
        "Timezone offset — autodeclarado + inferido por actividad → cruzar con IP",
        ("Un actor con buena opsec puede evadir 1–2 señales, pero raramente todas", 1),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SECCIÓN 3: LA INFRAESTRUCTURA
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "03", "La infraestructura",
                  "Por qué todo local y cómo montamos el entorno de análisis")

    content_slide(prs, "Por qué procesamos todo en local", [
        "Los datos contienen emails, contraseñas, IPs y contenido privado reales",
        "Enviar texto a APIs externas (OpenAI, Google, Cohere) está fuera de discusión",
        "GDPR + ética de investigación: minimizar exposición de datos de terceros",
        "Reproducibilidad: el análisis tiene que correr offline, sin dependencias externas",
        "Tenemos GPU (RTX A2000 12GB) → inferencia local es rápida y suficiente",
        ("Regla de oro: si no lo pondrías en un paper público, no lo mandes a una API", 1),
    ])

    two_col_slide(prs, "El stack de análisis", "Entorno", [
        "uv — gestor de entornos Python (10× más rápido que pip/venv)",
        "Python 3.12 — estabilidad + type hints modernos",
        "Jupyter Lab — análisis exploratorio, demo en vivo",
        "uv run jupyter — sin activar venv manualmente",
    ], "Librerías clave", [
        "pandas + numpy — manipulación y álgebra",
        "matplotlib + seaborn — visualización estática",
        "plotly — visualización interactiva (UMAP scatter)",
        "scikit-learn — clustering, métricas",
        "umap-learn — reducción dimensional 768D → 2D",
        "ollama — cliente para modelos locales",
    ])

    content_slide(prs, "Ollama: inferencia local de LLMs", [
        "Ollama = servidor local que expone una API compatible con OpenAI",
        "Descarga y gestiona modelos GGUF directamente: ollama pull nomic-embed-text",
        "nomic-embed-text — modelo de embeddings de 768 dimensiones",
        ("Entrenado específicamente para retrieval y clustering, no solo chat", 1),
        ("Alternativa: sentence-transformers — más flexible pero más lento en GPU", 1),
        "El análisis también funciona sin Ollama — estilometría tradicional como fallback",
        ("Ver src/stylometry.py — features manuales sin necesidad de GPU", 1),
    ], note="ollama pull nomic-embed-text (~300MB) · ollama pull qwen2.5:14b (~9GB)")

    content_slide(prs, "El parser: diseño y decisiones", [
        "Streaming puro — los dumps son 100–500MB, no caben en RAM completos",
        "Lee línea por línea, parsea VALUES(...) de cada INSERT",
        "El parser infiere el schema leyendo la primera fila — sin hardcodear columnas",
        "Auto-detect de encoding: peek 2KB → probar UTF-8, fallback cp1251",
        "Ignorar archivos ._filename.zip (artefactos de macOS en el dataset)",
        ("Bug resuelto: spaces después de coma en VALUES → el tokenizer confundía el estado", 1),
        ("Bug multiline: posts con \\n literal partían el parser → buffer + paren balancing", 1),
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SECCIÓN 4: TÉCNICAS DE ANÁLISIS
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "04", "Técnicas de análisis",
                  "De los datos en bruto a perfiles de actores")

    content_slide(prs, "Plan de análisis — 5 técnicas", [
        "1. Análisis de usuarios  —  demografía, actividad, datos de contacto",
        "2. Inferencia de timezone  —  actividad horaria → UTC offset → región probable",
        "3. Correlación cross-foro  —  mismo actor en múltiples foros via identifiers",
        "4. Estilometría tradicional  —  features de estilo extraídas manualmente",
        "5. Embeddings + UMAP  —  clustering de estilo con modelos de lenguaje locales",
        ("Las técnicas se complementan: ninguna sola es concluyente", 1),
        ("La convergencia de 2+ señales independientes es lo que construye un caso", 1),
    ])

    content_slide(prs, "1. Análisis de usuarios", [
        "Distribución temporal: cuándo se registraron → cuándo estuvo activo el foro",
        "Power law de posts: pocos usuarios muy activos, la gran mayoría inactivos",
        ("Los top-10 posters son los targets de mayor prioridad para análisis profundo", 1),
        "Contactos: cuántos dejaron email, ICQ, Skype, Jabber → densidad de señales",
        "Posts per user por foro → indica nivel de actividad real vs registros fantasma",
        "Visible vs hidden posts → posts moderados son señal de comportamiento problemático",
    ])

    content_slide(prs, "2. Inferencia de timezone", [
        "Premisa: la gente postea en horas activas — 08:00 a 23:00 hora local",
        "Para cada usuario: histograma de hora UTC de sus posts (de dateline)",
        "Buscamos el UTC offset que maximiza la actividad dentro de esa ventana de 15h",
        "Resultado: estimación probabilística del huso horario del usuario",
        ("No es determinístico — es una señal más, no una prueba", 1),
        "Validación: comparar con timezoneoffset autodeclarado en el perfil",
        ("Mismatch > 2h = señal de VPN o usuario que mintió en su perfil", 1),
    ])

    content_slide(prs, "3. Correlación cross-foro", [
        "Mismo actor registrado en múltiples foros = mayor credibilidad en la red criminal",
        "Señales por confiabilidad:",
        ("ICQ number — identificador duro, raramente cambiado → muy fiable", 1),
        ("Password hash — mismo hash en distintos foros → misma contraseña → mismo actor", 1),
        ("Email — throwaway pero si lo reusan es oro", 1),
        ("Username — útil solo con nicknames únicos y largos; colisión alta en nombres cortos", 1),
        "Resultado: grafo de identidades — nodos = usuarios, aristas = señales compartidas",
        ("Un nodo con alta centralidad es un actor relevante en múltiples comunidades", 1),
    ])

    content_slide(prs, "4. Estilometría", [
        "El estilo de escritura es difícil de cambiar consistentemente bajo presión",
        "Vocabulario, longitud de oraciones, uso de puntuación, errores tipográficos habituales",
        "Features manuales (src/stylometry.py):",
        ("Ratio puntuación/palabras, avg word length, function words, mayúsculas", 1),
        ("Operan sobre el texto tal cual — no necesitan GPU ni Ollama", 1),
        "Features semánticas (src/embeddings.py):",
        ("nomic-embed-text convierte todos los posts de un usuario en un vector de 768D", 1),
        ("Dos usuarios con cosine similarity > 0.85 tienen un estilo muy similar", 1),
    ])

    content_slide(prs, "5. UMAP — visualización de clusters", [
        "Problema: 768 dimensiones no se pueden visualizar directamente",
        "UMAP (Uniform Manifold Approximation and Projection) → reduce a 2D",
        "Preserva estructura local: puntos cercanos en 768D siguen cercanos en 2D",
        "Output: scatter plot interactivo (Plotly) — hover muestra username + foro",
        "Clusters visibles = grupos de usuarios con estilo de escritura similar",
        ("Un cluster con usuarios de varios foros distintos = posible red de actores", 1),
        ("Outliers = usuarios con estilo muy singular — interesantes para análisis manual", 1),
    ], note="Requiere Ollama con nomic-embed-text. Ver sección Setup del notebook.")

    # ══════════════════════════════════════════════════════════════════════════
    # SECCIÓN 5: DEMO EN VIVO
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "05", "Demo en vivo",
                  "Recorremos el notebook completo con los datos reales")

    demo_slide(prs, "01_carding_forums.ipynb — recorrido completo", [
        "Parsear los 10 dumps y verificar volumen de datos",
        "Explorar usuarios: demografía, contactos, top posters",
        "Inferir timezones y detectar mismatches (posibles VPN)",
        "Encontrar identidades cross-foro via ICQ, email y password hash",
        "Embeddings de estilo + reducción UMAP + scatter interactivo",
    ])

    # ══════════════════════════════════════════════════════════════════════════
    # SECCIÓN 6: HALLAZGOS Y ÉTICA
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, "06", "Hallazgos y ética",
                  "Qué encontramos, qué no podemos afirmar y cómo manejamos los datos")

    content_slide(prs, "Hallazgos — resumen cuantitativo", [
        "10 foros analizados · 1,138,681 usuarios · 1,527,311 posts",
        "Región dominante: UTC+3 → Rusia / Ucrania / Europa del Este",
        "Identidades cross-foro: [completar con resultados del notebook]",
        "Top señal de correlación: ICQ — más estable que email entre 2009 y 2016",
        "Post-2016: migraron a Telegram/Jabber — ICQ desaparece de los perfiles",
        "Patrón de actividad: picos en horario 10:00–20:00 UTC+3 (lunes a viernes)",
    ], note="Completar con cifras exactas tras correr el notebook end-to-end")

    content_slide(prs, "Limitaciones del análisis", [
        "Timezone inference: ±2h de error típico — VPN users generan ruido",
        "Username matching: alta colisión en nombres cortos (admin, user1, etc.)",
        "Embeddings: nomic-embed-text no está optimizado para texto cirílico",
        ("Los posts en ruso producen embeddings menos discriminativos", 1),
        "Carder.su: solo tabla user, 0 posts — no podemos hacer timezone ni estilometría",
        "CardingMafia / Cardmafia.cc: probable duplicado → puede inflar métricas cross-foro",
        "Dumps parciales: no sabemos qué tablas faltan en cada uno",
    ])

    content_slide(prs, "Marco legal y ético", [
        "Estos dumps son públicamente conocidos en el ámbito de ciberseguridad defensiva",
        "El análisis es defensivo: entender TTPs, identificar actores, proteger víctimas",
        "No publicar datos originales ni perfiles individuales sin justificación documentada",
        "Contexto de este material: formación interna — no investigación publicable",
        "Regla de mínima exposición: procesar solo lo necesario, borrar resultados sensibles post-análisis",
        ("Si en algún momento encontrás datos de víctimas reales — pará y consulta", 1),
    ])

    # ── Cierre ────────────────────────────────────────────────────────────────
    title_slide(prs,
        "Preguntas",
        "github.com/csbc26  ·  notebooks/cases/01_carding_forums.ipynb",
        "CSBC26  ·  OpHarvestSeason",
    )


# ─── Entry point ──────────────────────────────────────────────────────────────

def main():
    template_path = sys.argv[1] if len(sys.argv) > 1 else None

    if template_path:
        prs = Presentation(template_path)
        # Preserve template slide size
    else:
        prs = Presentation()
        prs.slide_width  = W
        prs.slide_height = H

    build(prs)

    out = Path(__file__).parent / "csbc26_carding_forums.pptx"
    prs.save(str(out))
    print(f"Saved: {out}  ({out.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
