"""
Tema visual compartido para todas las presentaciones CSBC26.

Construye sobre la plantilla oficial del Cybersecurity Summer Bootcamp 2026
(_shared/plantilla_csbc26.pptx): el fondo de marca, el logo, los patrocinadores
y los layouts ya vienen del máster de PowerPoint. Este módulo solo puebla los
marcadores de posición de cada layout y añade los elementos que la plantilla
no cubre de forma nativa (tablas, bloques de código, dos columnas).

Uso:
    from _shared.theme import *
    prs = new_presentation()
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE_PATH = Path(__file__).parent / "plantilla_csbc26.pptx"

# ─── Paleta de colores ────────────────────────────────────────────────────────
# Valores extraídos de los propios layouts de la plantilla (lstStyle de cada
# marcador de posición), no inventados: así el texto que añadimos a mano
# (tablas, código, dos columnas) combina con lo que la plantilla ya trae.
# Diapositivas oscuras (portada, sección, cierre) — fondo de marca ya incluido.
C_TITLE  = RGBColor(0xFF, 0xFF, 0xFF)  # blanco — títulos sobre fondo oscuro
C_BODY   = RGBColor(0xD6, 0xDA, 0xE2)  # gris claro — cuerpo sobre fondo oscuro
C_MUTED  = RGBColor(0x9A, 0xA3, 0xB5)  # atenuado sobre fondo oscuro

# Diapositivas claras (contenido, tablas, código) — fondo blanco/gris de marca.
C_TITLE_D = RGBColor(0xCD, 0x2D, 0x37)  # rojo de marca — título del layout "Bullets"
C_BODY_D  = RGBColor(0x46, 0x54, 0x61)  # gris pizarra — cuerpo del layout "Bullets"
C_MUTED_D = RGBColor(0x70, 0x78, 0x88)  # atenuado — sobre fondo claro

C_ACCENT  = RGBColor(0xCD, 0x2D, 0x37)  # mismo rojo de marca — destacados y tablas
C_CODE_BG = RGBColor(0x0D, 0x1B, 0x2A)  # bloque de código, sobre cualquier fondo

FONT = "Arial"

# ─── Dimensiones de diapositiva (16:9 panorámica, iguales a la plantilla) ────
W = Inches(13.33)
H = Inches(7.5)

# ─── Layouts de la plantilla usados por este módulo ──────────────────────────
LAYOUT_TITLE   = "Portada presentación"
LAYOUT_SECTION = "Portada sección 1"
LAYOUT_CONTENT = "Bullets"


def new_presentation() -> Presentation:
    """Crea una presentación nueva a partir de la plantilla de marca, sin las diapositivas de ejemplo que trae."""
    prs = Presentation(str(TEMPLATE_PATH))
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        prs.part.drop_rel(sldId.rId)
        xml_slides.remove(sldId)
    return prs


def _layout(prs: Presentation, name: str):
    """Busca un layout de la plantilla por nombre."""
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout {name!r} no encontrado en la plantilla")


def _set_notes(slide, notes: str) -> None:
    """Escribe el guion de oratoria en las notas del orador (modo presentación)."""
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ─── Auxiliares de bajo nivel ─────────────────────────────────────────────────

def _add_rect(slide, x, y, w, h, fill: RGBColor):
    """Agrega un rectángulo relleno sin borde."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def _add_txbox(
    slide,
    text: str,
    x, y, w, h,
    size: int = 20,
    color: RGBColor = C_BODY_D,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    """Agrega un cuadro de texto de un solo párrafo."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def _set_title_placeholder(slide, idx: int, text: str) -> None:
    """
    Escribe el título en el marcador de posición idx de un layout claro.
    No fija fuente/tamaño/color a propósito: el layout ya define Arial 40pt
    negrita en rojo de marca (CD2D37) para este marcador — lo heredamos.
    """
    ph = slide.placeholders[idx]
    ph.text_frame.word_wrap = True
    run = ph.text_frame.paragraphs[0].add_run()
    run.text = text


# El marcador de título del layout mide ~7.46in de ancho y 0.57in de alto,
# pensado para un título corto de una línea a 40pt. Con nuestros títulos más
# largos (heredados del diseño anterior) puede envolver a una segunda o tercera
# línea y chocar con el contenido de debajo si no le hacemos hueco.
_TITLE_CHARS_PER_LINE = 26
_TITLE_LINE_PUSH_IN = 0.55
# Aire mínimo entre el título y el cuerpo aunque el título quepa en una sola
# línea — sin esto, el contenido queda pegado justo debajo del título.
_TITLE_BASE_GAP_IN = 0.15


def _title_push_in(title: str) -> float:
    """Cuánto hay que bajar el cuerpo de una diapositiva 'Bullets' para dejar
    aire tras el título, sume o no líneas de más."""
    return _TITLE_BASE_GAP_IN + _TITLE_LINE_PUSH_IN * _title_extra_lines(title)


def _title_extra_lines(text: str) -> int:
    """
    Estima cuántas líneas de más ocupará el título a 40pt en su marcador.
    Simula el ajuste de línea palabra por palabra (no basta con dividir la
    longitud total entre el ancho: una palabra larga puede forzar un salto
    antes de llenar la línea, como en "UMAP + HDBSCAN — estructura...").
    """
    lines, cur = 1, 0
    for word in text.split():
        wl = len(word)
        if cur == 0:
            cur = wl
        elif cur + 1 + wl <= _TITLE_CHARS_PER_LINE:
            cur += 1 + wl
        else:
            lines += 1
            cur = wl
    return lines - 1


def _add_bullets(
    slide,
    ph_idx: int,
    bullets: list,
    size: int = 16,
) -> None:
    """
    Puebla un marcador de posición de cuerpo con una lista de viñetas.
    bullets puede contener str (normal) o ("texto", nivel_sangría) para sub-viñetas.

    Nivel 0 hereda el estilo nativo del layout (Arial 18pt, gris pizarra,
    viñeta-imagen del rombo rojo de marca) — no tocamos fuente/tamaño/color
    ni desactivamos la viñeta. Los sub-niveles no tienen estilo propio en el
    layout, así que ahí sí aplicamos un prefijo y un tamaño más discretos.
    """
    tf = slide.placeholders[ph_idx].text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if level == 0:
            run = p.add_run()
            run.text = text
            continue

        p.level = level
        pPr = p._p.get_or_add_pPr()
        indent = Emu(Inches(0.3 * level))
        pPr.set("marL", str(int(indent)))
        pPr.set("indent", str(int(-Inches(0.25))))
        # Los sub-niveles no tienen viñeta-imagen propia en el layout (caerían
        # en el bullet de texto negro del máster); la desactivamos y ponemos
        # un prefijo tipográfico discreto en su lugar.
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))

        run = p.add_run()
        run.text = "·  " + text
        run.font.size = Pt(size - level)
        run.font.color.rgb = C_MUTED_D
        run.font.name = FONT
        run.font.bold = False


# ─── Constructores de diapositivas ────────────────────────────────────────────

def title_slide(prs: Presentation, title: str, subtitle: str, tag: str, notes: str = "") -> None:
    """
    Diapositiva de portada — usa el layout de marca 'Portada presentación'.
    El título hereda el estilo nativo del layout (Arial 48pt negrita blanco,
    alineado a la izquierda); subtítulo y tag son líneas añadidas por nosotros
    y sí llevan tamaño/color explícitos, más pequeños y discretos.
    """
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_TITLE))
    tf = slide.placeholders[10].text_frame
    tf.word_wrap = True

    r = tf.paragraphs[0].add_run()
    r.text = title

    p1 = tf.add_paragraph()
    r = p1.add_run()
    r.text = subtitle
    r.font.size = Pt(18)
    r.font.color.rgb = C_BODY
    r.font.name = FONT

    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = tag
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = C_MUTED
    r.font.name = FONT

    _set_notes(slide, notes)


def section_slide(prs: Presentation, number: str, title: str, subtitle: str = "", notes: str = "") -> None:
    """
    Diapositiva divisora de sección — usa el layout de marca 'Portada sección 1'.
    El título hereda el estilo nativo (Arial 36pt negrita blanco, centrado);
    la etiqueta "SECCIÓN N" y el subtítulo son líneas propias, centradas para
    acompañar la alineación que define el layout.
    """
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_SECTION))
    tf = slide.placeholders[10].text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = f"SECCIÓN {number}"
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = C_ACCENT
    r.font.name = FONT

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r = p1.add_run()
    r.text = title

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run()
        r.text = subtitle
        r.font.size = Pt(16)
        r.font.color.rgb = C_BODY
        r.font.name = FONT

    _set_notes(slide, notes)


def content_slide(
    prs: Presentation,
    title: str,
    bullets: list,
    note: str = "",
    notes: str = "",
) -> None:
    """Diapositiva de contenido — layout 'Bullets': título + lista de viñetas + nota opcional al pie."""
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_CONTENT))
    _set_title_placeholder(slide, 13, title)

    push = Inches(_title_push_in(title))
    body_ph = slide.placeholders[16]
    # Hay que reafirmar left/width explícitamente: si solo se asigna
    # top/height, python-pptx crea un xfrm nuevo con left/width en 0
    # (el texto queda en una caja de ancho cero, invisible al exportar).
    left, width = body_ph.left, body_ph.width
    body_ph.top = body_ph.top + push
    body_ph.height = body_ph.height - push
    body_ph.left = left
    body_ph.width = width

    _add_bullets(slide, 16, bullets)
    if note:
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(1.15), Inches(6.3), Inches(11.0), Inches(0.4),
                   size=11, color=C_MUTED_D)
    _set_notes(slide, notes)


def two_col_slide(
    prs: Presentation,
    title: str,
    left_title: str, left_items: list[str],
    right_title: str, right_items: list[str],
    note: str = "",
    notes: str = "",
) -> None:
    """Diapositiva de dos columnas — layout 'Bullets' con título global y dos listas independientes."""
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_CONTENT))
    _set_title_placeholder(slide, 13, title)
    push_in = _title_push_in(title)

    col_w = Inches(5.3)
    col_h = Inches(4.2 - push_in)

    for col_x, col_title, items in [
        (Inches(1.15), left_title,  left_items),
        (Inches(6.8),  right_title, right_items),
    ]:
        _add_txbox(slide, col_title, col_x, Inches(1.85 + push_in), col_w, Inches(0.45),
                   size=16, color=C_ACCENT, bold=True)

        txb = slide.shapes.add_textbox(col_x, Inches(2.35 + push_in), col_w, col_h)
        tf = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Prefijo en el rojo de marca, texto en el color de cuerpo — mismo
            # tratamiento que la viñeta nativa de content_slide, para que no
            # se sienta como un icono distinto entre diapositivas.
            bullet_run = p.add_run()
            bullet_run.text = "▸  "
            bullet_run.font.size = Pt(16)
            bullet_run.font.color.rgb = C_ACCENT
            bullet_run.font.name = FONT
            text_run = p.add_run()
            text_run.text = item
            text_run.font.size = Pt(16)
            text_run.font.color.rgb = C_BODY_D
            text_run.font.name = FONT

    if note:
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(1.15), Inches(6.3), Inches(11.0), Inches(0.4),
                   size=11, color=C_MUTED_D)

    _set_notes(slide, notes)


def _style_table_cell(cell, text: str, bg: RGBColor, color: RGBColor, size: int, bold: bool) -> None:
    """Rellena y da formato a una celda de una tabla nativa de PowerPoint."""
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.margin_left = Emu(45000)
    cell.margin_right = Emu(45000)
    cell.margin_top = Emu(25000)
    cell.margin_bottom = Emu(25000)
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    # Una celda vacía no genera runs al asignar cell.text — sin este caso
    # especial, aplicar el formato revienta con IndexError.
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    note: str = "",
    notes: str = "",
) -> None:
    """
    Diapositiva con tabla de datos — layout 'Bullets' con encabezado destacado
    y filas alternadas. Usa una tabla nativa de PowerPoint (no formas dibujadas
    a mano) para que columnas y filas se puedan ajustar directamente en
    PowerPoint/LibreOffice; el ajuste de línea y el alto de fila los resuelve
    la propia aplicación.
    """
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_CONTENT))
    _set_title_placeholder(slide, 13, title)
    push_in = _title_push_in(title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h = Inches(0.4)
    table_top = Inches(1.8 + push_in)

    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, Inches(1.15), table_top, Inches(11.0), row_h * n_rows,
    )
    table = graphic_frame.table
    # Desactivamos el estilo de tabla por defecto (colores del tema, bandas
    # automáticas) para controlar nosotros el color de cada celda.
    table.first_row = False
    table.horz_banding = False

    col_w = Inches(11.0 / n_cols)
    for col in table.columns:
        col.width = col_w
    for row in table.rows:
        row.height = row_h

    for j, h in enumerate(headers):
        _style_table_cell(table.cell(0, j), h, C_ACCENT, C_TITLE, 12, True)

    for i, row_cells in enumerate(rows):
        bg = RGBColor(0xF1, 0xF2, 0xF6) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for j, cell_text in enumerate(row_cells):
            _style_table_cell(table.cell(i + 1, j), cell_text, bg, C_BODY_D, 11, False)

    if note:
        note_top_in = 1.8 + push_in + 0.4 * n_rows + 0.15
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(1.15), Inches(max(6.3, note_top_in)), Inches(11.0), Inches(0.35),
                   size=10, color=C_MUTED_D)

    _set_notes(slide, notes)


# Paleta estándar de tier list (rojo = top tier, azul/gris = tier más bajo),
# la misma convención que usan las tier-list maker habituales.
TIER_COLORS = {
    "S": RGBColor(0xFF, 0x7F, 0x7F),
    "A": RGBColor(0xFF, 0xBF, 0x7F),
    "B": RGBColor(0xFF, 0xDF, 0x7F),
    "C": RGBColor(0xFF, 0xFF, 0x7F),
    "D": RGBColor(0xBF, 0xFF, 0x7F),
    "F": RGBColor(0xBF, 0xBF, 0xBF),
}


def tierlist_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    note: str = "",
    notes: str = "",
) -> None:
    """
    Diapositiva de tier list — misma tabla nativa que table_slide, pero cada
    fila se colorea según el tier de su primera columna (S/A/B/C/D/F), al
    estilo de las tier-list maker habituales. La columna de tier va más ancha
    y centrada para que la letra destaque como una insignia.
    """
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_CONTENT))
    _set_title_placeholder(slide, 13, title)
    push_in = _title_push_in(title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h = Inches(0.4)
    table_top = Inches(1.8 + push_in)

    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, Inches(1.15), table_top, Inches(11.0), row_h * n_rows,
    )
    table = graphic_frame.table
    table.first_row = False
    table.horz_banding = False

    tier_col_w = Inches(0.9)
    other_col_w = Emu(int((Inches(11.0) - tier_col_w) / (n_cols - 1)))
    for j, col in enumerate(table.columns):
        col.width = tier_col_w if j == 0 else other_col_w
    for row in table.rows:
        row.height = row_h

    for j, h in enumerate(headers):
        _style_table_cell(table.cell(0, j), h, C_ACCENT, C_TITLE, 12, True)
        table.cell(0, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT

    for i, row_cells in enumerate(rows):
        tier = row_cells[0].strip().upper()
        bg = TIER_COLORS.get(tier, RGBColor(0xF1, 0xF2, 0xF6))
        for j, cell_text in enumerate(row_cells):
            is_tier_badge = j == 0
            _style_table_cell(
                table.cell(i + 1, j), cell_text, bg, C_BODY_D,
                18 if is_tier_badge else 11,
                bold=is_tier_badge,
            )
            if is_tier_badge:
                table.cell(i + 1, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if note:
        note_top_in = 1.8 + push_in + 0.4 * n_rows + 0.15
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(1.15), Inches(max(6.3, note_top_in)), Inches(11.0), Inches(0.35),
                   size=10, color=C_MUTED_D)

    _set_notes(slide, notes)


def code_slide(
    prs: Presentation,
    title: str,
    description: str,
    code_lines: list[str],
    note: str = "",
    notes: str = "",
) -> None:
    """Diapositiva con bloque de código monoespacio — layout 'Bullets'."""
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_CONTENT))
    _set_title_placeholder(slide, 13, title)
    push_in = _title_push_in(title)

    y = Inches(1.85 + push_in)
    if description:
        _add_txbox(slide, description, Inches(1.15), y, Inches(11.0), Inches(0.55),
                   size=16, color=C_BODY_D)
        y = y + Inches(0.6)

    code_h = Inches(6.35) - y

    _add_rect(slide, Inches(1.1), y, Inches(11.1), code_h, C_CODE_BG)

    code_text = "\n".join(code_lines)
    txb = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.15),
                                   Inches(10.7), code_h - Inches(0.3))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)
    run.font.name = "Consolas"

    if note:
        _add_txbox(slide, f"ℹ  {note}",
                   Inches(1.15), Inches(6.3), Inches(11.0), Inches(0.4),
                   size=11, color=C_MUTED_D)

    _set_notes(slide, notes)


def demo_slide(prs: Presentation, section_title: str, steps: list[str], notebook_path: str = "", notes: str = "") -> None:
    """
    Diapositiva de demo en vivo — reutiliza el layout oscuro 'Portada sección 1'.
    El título hereda el estilo nativo (Arial 36pt negrita blanco, centrado).
    """
    slide = prs.slides.add_slide(_layout(prs, LAYOUT_SECTION))
    tf = slide.placeholders[10].text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = "DEMO EN VIVO"
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = C_ACCENT
    r.font.name = FONT

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r = p1.add_run()
    r.text = section_title

    if notebook_path:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run()
        r.text = notebook_path
        r.font.size = Pt(14)
        r.font.color.rgb = C_MUTED
        r.font.name = FONT

    txb = slide.shapes.add_textbox(Inches(2.74), Inches(4.0), Inches(7.84), Inches(2.8))
    tf2 = txb.text_frame
    tf2.word_wrap = True
    for i, step in enumerate(steps):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run = p.add_run()
        run.text = f"{i+1}.  {step}"
        run.font.size = Pt(16)
        run.font.color.rgb = C_BODY
        run.font.name = FONT

    _set_notes(slide, notes)
